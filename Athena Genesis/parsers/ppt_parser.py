# -*- coding: utf-8 -*-
"""
PowerPoint文件解析器
"""

import os
from typing import List, Dict, Any
from .base_parser import BaseParser
from config.settings import SETTINGS


class PPTParser(BaseParser):
    """PowerPoint文件解析器"""

    def __init__(self):
        super().__init__()
        self.supported_extensions = SETTINGS.SUPPORTED_EXTENSIONS['ppt']
        self.parser_name = "PPTParser"
        self.has_pptx = False

        self._check_dependencies()

    def _check_dependencies(self) -> None:
        """检查依赖"""
        try:
            import pptx
            self.has_pptx = True
        except ImportError:
            self.has_pptx = False

    def can_parse(self, file_path: str) -> bool:
        """检查是否能解析PPT文件"""
        if not self.has_pptx:
            return False

        ext = os.path.splitext(file_path)[1].lower()
        return ext in ['.pptx', '.ppt']

    def parse(self, file_path: str) -> str:
        """解析PowerPoint文件"""
        if not self.has_pptx:
            return "[需要安装python-pptx库才能解析PPT文件]"

        try:
            from pptx import Presentation

            prs = Presentation(file_path)

            content_lines = [f"[PowerPoint演示文稿: {os.path.basename(file_path)}]"]
            content_lines.append(f"幻灯片数量: {len(prs.slides)}")

            # 解析幻灯片
            for i, slide in enumerate(prs.slides[:20]):  # 最多20张幻灯片
                slide_content = self._parse_slide(slide, i + 1)
                if slide_content:
                    content_lines.extend(slide_content)

            if len(prs.slides) > 20:
                remaining = len(prs.slides) - 20
                content_lines.append(f"\n...省略 {remaining} 张幻灯片")

            return "\n".join(content_lines)

        except Exception as e:
            return f"[PPT解析错误]: {str(e)}"

    def _parse_slide(self, slide, slide_num: int) -> List[str]:
        """解析单张幻灯片"""
        content = []
        content.append(f"\n--- 幻灯片 {slide_num} ---")

        # 幻灯片标题
        if hasattr(slide, 'shapes'):
            title_found = False
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()

                    # 检查是否是标题（通常是第一个文本框或占位符）
                    if not title_found and len(text) < 100:
                        content.append(f"标题: {text}")
                        title_found = True
                    else:
                        # 限制文本长度
                        if len(text) > 200:
                            text = text[:200] + "..."
                        content.append(text)

        # 如果没有任何内容
        if len(content) == 1:
            content.append("[无文本内容]")

        return content

    def get_presentation_summary(self, file_path: str) -> Dict[str, Any]:
        """获取演示文稿摘要"""
        if not self.has_pptx:
            return {"error": "需要安装python-pptx库"}

        try:
            from pptx import Presentation

            prs = Presentation(file_path)

            summary = {
                "filename": os.path.basename(file_path),
                "slide_count": len(prs.slides),
                "slides": [],
                "layout_types": set(),
                "has_notes": False
            }

            # 分析幻灯片
            for i, slide in enumerate(prs.slides[:10]):  # 最多10张
                slide_info = {
                    "slide_number": i + 1,
                    "layout": slide.slide_layout.name,
                    "shapes_count": len(slide.shapes),
                    "has_notes": hasattr(slide, 'notes_slide') and slide.notes_slide is not None
                }

                summary["layout_types"].add(slide.slide_layout.name)
                summary["slides"].append(slide_info)

                if slide_info["has_notes"]:
                    summary["has_notes"] = True

            summary["layout_types"] = list(summary["layout_types"])

            return summary

        except Exception as e:
            return {"error": str(e)}